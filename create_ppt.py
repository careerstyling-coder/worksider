#!/usr/bin/env python3
"""Workside 서비스 소개 PPT 생성 스크립트"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Colors
BLUE = RGBColor(0x18, 0x77, 0xF2)
DARK = RGBColor(0x1A, 0x1A, 0x2E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xF7, 0xF8, 0xFA)
GREEN = RGBColor(0x31, 0xA2, 0x4C)
YELLOW = RGBColor(0xF7, 0xB9, 0x28)
PURPLE = RGBColor(0x9B, 0x59, 0xB6)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

def add_bg(slide, color=WHITE):
    """Add background color to slide"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_bg(slide, color):
    """Add a full-slide shape as background"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False, color=DARK, alignment=PP_ALIGN.LEFT, font_name="Arial"):
    """Add a text box with formatting"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_accent_bar(slide, left, top, width=Inches(0.6), height=Inches(0.06)):
    """Add blue accent bar"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = BLUE
    shape.line.fill.background()

def add_card(slide, left, top, width, height, color=WHITE):
    """Add a rounded rectangle card"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
    shape.line.width = Pt(1)
    shape.shadow.inherit = False
    return shape

# ═══════════════════════════════════════════════════
# SLIDE 1: Title
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_shape_bg(slide, WHITE)

# Blue accent top bar
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.08))
shape.fill.solid()
shape.fill.fore_color.rgb = BLUE
shape.line.fill.background()

add_text_box(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(0.6),
             "Workside", 20, bold=True, color=BLUE)

add_text_box(slide, Inches(1.5), Inches(2.5), Inches(10), Inches(1.5),
             "나를 발견하는 첫 걸음,\nWorkstyle DNA", 48, bold=True, color=DARK)

add_text_box(slide, Inches(1.5), Inches(4.3), Inches(8), Inches(0.8),
             "일을 더욱 일답게 하고자 하는 직장인들의 성장 커뮤니티 플랫폼", 22, color=GRAY)

add_accent_bar(slide, Inches(1.5), Inches(5.5))

add_text_box(slide, Inches(1.5), Inches(5.8), Inches(6), Inches(0.5),
             "2026.03  |  Service Introduction", 14, color=GRAY)

# ═══════════════════════════════════════════════════
# SLIDE 2: Problem
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape_bg(slide, WHITE)

add_accent_bar(slide, Inches(1), Inches(0.8))
add_text_box(slide, Inches(1), Inches(1.0), Inches(10), Inches(0.6),
             "Problem", 14, bold=True, color=BLUE)
add_text_box(slide, Inches(1), Inches(1.4), Inches(10), Inches(0.8),
             "왜 우리는 일할 때 서로를 이해하기 어려울까?", 36, bold=True, color=DARK)

# Problem cards
problems = [
    ("나를 모른다", "내가 일을 어떻게 하는 사람인지\n설명할 언어가 없다"),
    ("다름을 모른다", "동료와의 업무 스타일 차이를\n갈등으로 오해한다"),
    ("성장 방향을 모른다", "나에게 맞는 성장 방법이\n무엇인지 알기 어렵다"),
]

for i, (title, desc) in enumerate(problems):
    left = Inches(1 + i * 3.8)
    card = add_card(slide, left, Inches(2.8), Inches(3.4), Inches(3.2))

    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.3), Inches(3.1), Inches(0.5), Inches(0.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = BLUE
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].text = str(i + 1)
    tf.paragraphs[0].font.size = Pt(16)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    add_text_box(slide, left + Inches(0.3), Inches(3.8), Inches(2.8), Inches(0.5),
                 title, 22, bold=True, color=DARK)
    add_text_box(slide, left + Inches(0.3), Inches(4.4), Inches(2.8), Inches(1.2),
                 desc, 15, color=GRAY)

# ═══════════════════════════════════════════════════
# SLIDE 3: Solution - Workstyle DNA
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape_bg(slide, WHITE)

add_accent_bar(slide, Inches(1), Inches(0.8))
add_text_box(slide, Inches(1), Inches(1.0), Inches(10), Inches(0.6),
             "Solution", 14, bold=True, color=BLUE)
add_text_box(slide, Inches(1), Inches(1.4), Inches(10), Inches(0.8),
             "Workstyle DNA로 나를 발견하다", 36, bold=True, color=DARK)

add_text_box(slide, Inches(1), Inches(2.4), Inches(10), Inches(0.6),
             "4가지 축으로 나의 업무 성향을 과학적으로 분석합니다", 18, color=GRAY)

# 4 axes
axes = [
    ("P", "실행력", "목표 달성과 결과물에\n집중하는 성향", BLUE),
    ("C", "협력", "동료와의 소통과\n팀워크 지향성", GREEN),
    ("Pol", "영향력", "조직 역학을 읽고\n관계를 활용하는 능력", YELLOW),
    ("S", "자율성", "독립적으로 일하며\n자기 주도적 성향", PURPLE),
]

for i, (code, name, desc, color) in enumerate(axes):
    left = Inches(1 + i * 3)

    # Icon circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.8), Inches(3.3), Inches(0.7), Inches(0.7))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].text = code
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    add_text_box(slide, left, Inches(4.2), Inches(2.6), Inches(0.5),
                 name, 22, bold=True, color=DARK, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left, Inches(4.7), Inches(2.6), Inches(1.0),
                 desc, 14, color=GRAY, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════
# SLIDE 4: 6 Personas
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape_bg(slide, WHITE)

add_accent_bar(slide, Inches(1), Inches(0.8))
add_text_box(slide, Inches(1), Inches(1.0), Inches(10), Inches(0.6),
             "Personas", 14, bold=True, color=BLUE)
add_text_box(slide, Inches(1), Inches(1.4), Inches(10), Inches(0.8),
             "6가지 업무 페르소나", 36, bold=True, color=DARK)

personas = [
    ("전략적 성과자", "목표를 향해 달리는", BLUE),
    ("실무형 전문가", "깊이로 승부하는", RGBColor(0x06, 0xB6, 0xD4)),
    ("협력적 조정자", "함께 만들어가는", GREEN),
    ("자율형 독립가", "내 길을 가는", PURPLE),
    ("조직형 정치인", "흐름을 읽는", YELLOW),
    ("중도형 균형가", "균형을 찾는", RGBColor(0xEC, 0x48, 0x99)),
]

for i, (name, desc, color) in enumerate(personas):
    col = i % 3
    row = i // 3
    left = Inches(1 + col * 3.8)
    top = Inches(2.6 + row * 2.2)

    card = add_card(slide, left, top, Inches(3.4), Inches(1.8))

    # Color bar on left
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left + Inches(0.15), top + Inches(0.3), Inches(0.08), Inches(1.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()

    add_text_box(slide, left + Inches(0.5), top + Inches(0.4), Inches(2.6), Inches(0.5),
                 name, 20, bold=True, color=DARK)
    add_text_box(slide, left + Inches(0.5), top + Inches(1.0), Inches(2.6), Inches(0.5),
                 desc, 15, color=GRAY)

# ═══════════════════════════════════════════════════
# SLIDE 5: How It Works
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape_bg(slide, WHITE)

add_accent_bar(slide, Inches(1), Inches(0.8))
add_text_box(slide, Inches(1), Inches(1.0), Inches(10), Inches(0.6),
             "How It Works", 14, bold=True, color=BLUE)
add_text_box(slide, Inches(1), Inches(1.4), Inches(10), Inches(0.8),
             "성장의 선순환", 36, bold=True, color=DARK)

steps = [
    ("1", "발견", "DNA 진단으로\n나를 이해해요", "3분 세미 진단\n12문항 리커트 척도"),
    ("2", "공유", "동료와 서로를\n이해해요", "결과 링크 공유\n유형 비교"),
    ("3", "질문", "함께 고민하고\n답을 찾아요", "주간 질문 참여\n페르소나별 인사이트"),
    ("4", "성장", "매번 새로운\n나를 발견해요", "재진단으로 변화 추적\n커리어 성장"),
]

for i, (num, title, desc, detail) in enumerate(steps):
    left = Inches(0.8 + i * 3.1)

    # Step circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.9), Inches(2.8), Inches(0.8), Inches(0.8))
    circle.fill.solid()
    circle.fill.fore_color.rgb = BLUE
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Arrow (except last)
    if i < 3:
        add_text_box(slide, left + Inches(2.5), Inches(2.9), Inches(0.5), Inches(0.6),
                     "→", 28, color=RGBColor(0xCC, 0xCC, 0xCC), alignment=PP_ALIGN.CENTER)

    add_text_box(slide, left, Inches(3.8), Inches(2.6), Inches(0.5),
                 title, 24, bold=True, color=DARK, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left, Inches(4.3), Inches(2.6), Inches(0.8),
                 desc, 15, color=GRAY, alignment=PP_ALIGN.CENTER)

    # Detail card
    card = add_card(slide, left + Inches(0.1), Inches(5.3), Inches(2.4), Inches(1.2), LIGHT_GRAY)
    add_text_box(slide, left + Inches(0.3), Inches(5.5), Inches(2.0), Inches(0.8),
                 detail, 12, color=GRAY, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════
# SLIDE 6: Key Features
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape_bg(slide, WHITE)

add_accent_bar(slide, Inches(1), Inches(0.8))
add_text_box(slide, Inches(1), Inches(1.0), Inches(10), Inches(0.6),
             "Key Features", 14, bold=True, color=BLUE)
add_text_box(slide, Inches(1), Inches(1.4), Inches(10), Inches(0.8),
             "주요 기능", 36, bold=True, color=DARK)

features = [
    ("Workstyle DNA 진단", "3분 세미 / 10분 풀 진단\n4축 분석 + 페르소나 판정\n레이더 차트 시각화"),
    ("질문 피드", "주간 질문으로 다양한 시각 탐색\n페르소나별 응답 인사이트\n궁금합니다(제안) + Shout out"),
    ("결과 공유", "고유 링크로 결과 공유\n동료와 업무 스타일 비교\n팀 내 상호 이해 촉진"),
    ("관리자 대시보드", "실시간 통계 모니터링\n질문 생성/배포/마감 관리\n사용자 분석 (산업/규모별)"),
]

for i, (title, desc) in enumerate(features):
    col = i % 2
    row = i // 2
    left = Inches(1 + col * 5.8)
    top = Inches(2.6 + row * 2.4)

    card = add_card(slide, left, top, Inches(5.4), Inches(2.0))

    add_text_box(slide, left + Inches(0.4), top + Inches(0.3), Inches(4.6), Inches(0.5),
                 title, 20, bold=True, color=DARK)

    # Blue dot
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.2), top + Inches(0.4), Inches(0.12), Inches(0.12))
    dot.fill.solid()
    dot.fill.fore_color.rgb = BLUE
    dot.line.fill.background()

    add_text_box(slide, left + Inches(0.4), top + Inches(0.8), Inches(4.6), Inches(1.0),
                 desc, 14, color=GRAY)

# ═══════════════════════════════════════════════════
# SLIDE 7: Screenshots
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape_bg(slide, LIGHT_GRAY)

add_accent_bar(slide, Inches(1), Inches(0.8))
add_text_box(slide, Inches(1), Inches(1.0), Inches(10), Inches(0.6),
             "Product", 14, bold=True, color=BLUE)
add_text_box(slide, Inches(1), Inches(1.4), Inches(10), Inches(0.8),
             "서비스 화면", 36, bold=True, color=DARK)

# Landing screenshot
img_path = "/Users/maegbugpeulo/worksider/screenshots/01-landing-page.jpg"
if os.path.exists(img_path):
    slide.shapes.add_picture(img_path, Inches(0.8), Inches(2.5), Inches(5.8))
    add_text_box(slide, Inches(0.8), Inches(6.5), Inches(5.8), Inches(0.4),
                 "Landing Page", 12, color=GRAY, alignment=PP_ALIGN.CENTER)

# About screenshot
img_path2 = "/Users/maegbugpeulo/worksider/screenshots/02-about-page.jpg"
if os.path.exists(img_path2):
    slide.shapes.add_picture(img_path2, Inches(7), Inches(2.5), Inches(5.8))
    add_text_box(slide, Inches(7), Inches(6.5), Inches(5.8), Inches(0.4),
                 "About Page", 12, color=GRAY, alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════
# SLIDE 8: Tech Stack
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape_bg(slide, WHITE)

add_accent_bar(slide, Inches(1), Inches(0.8))
add_text_box(slide, Inches(1), Inches(1.0), Inches(10), Inches(0.6),
             "Technology", 14, bold=True, color=BLUE)
add_text_box(slide, Inches(1), Inches(1.4), Inches(10), Inches(0.8),
             "기술 스택 & 인프라", 36, bold=True, color=DARK)

techs = [
    ("Frontend", "Next.js 16 (App Router)\nReact 19\nTailwindCSS 4\nRecharts"),
    ("Backend", "Next.js API Routes\nZod Validation\nSupabase Auth (JWT)\nRow Level Security"),
    ("Database", "PostgreSQL 17\n11 Tables + RLS\nDB Migrations\nReal-time Subscriptions"),
    ("Infrastructure", "Vercel (Hosting)\nSupabase Cloud (DB)\nAuto SSL\nCI/CD (git push)"),
]

for i, (title, desc) in enumerate(techs):
    left = Inches(0.8 + i * 3.1)

    card = add_card(slide, left, Inches(2.6), Inches(2.8), Inches(4.0))

    # Header bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(2.6), Inches(2.8), Inches(0.6))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE
    bar.line.fill.background()

    add_text_box(slide, left + Inches(0.3), Inches(2.7), Inches(2.2), Inches(0.4),
                 title, 16, bold=True, color=WHITE)
    add_text_box(slide, left + Inches(0.3), Inches(3.5), Inches(2.2), Inches(2.8),
                 desc, 14, color=GRAY)

# ═══════════════════════════════════════════════════
# SLIDE 9: Growth Roadmap
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape_bg(slide, WHITE)

add_accent_bar(slide, Inches(1), Inches(0.8))
add_text_box(slide, Inches(1), Inches(1.0), Inches(10), Inches(0.6),
             "Roadmap", 14, bold=True, color=BLUE)
add_text_box(slide, Inches(1), Inches(1.4), Inches(10), Inches(0.8),
             "서비스 확장 로드맵", 36, bold=True, color=DARK)

phases = [
    ("Phase 1\n현재", "커뮤니티 MVP", "DNA 진단 + 질문 피드\n제안 시스템\n관리자 대시보드", BLUE),
    ("Phase 2\n3~6개월", "프로필 & 유료화", "실명 인증 시스템\n포트폴리오 작성/열람\n프리미엄 구독", GREEN),
    ("Phase 3\n6~12개월", "회사 탐색", "조직문화 리포트\n회사 탐색 (이직 준비)\n참여자 수익 공유", YELLOW),
    ("Phase 4\n12개월+", "채용 연결", "스카웃 제의\n면접 요청\n매칭 알고리즘", PURPLE),
]

for i, (phase, title, desc, color) in enumerate(phases):
    left = Inches(0.6 + i * 3.15)

    # Phase card
    card = add_card(slide, left, Inches(2.6), Inches(2.9), Inches(4.2))

    # Color top bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(2.6), Inches(2.9), Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()

    add_text_box(slide, left + Inches(0.3), Inches(2.9), Inches(2.3), Inches(0.7),
                 phase, 12, color=color)
    add_text_box(slide, left + Inches(0.3), Inches(3.6), Inches(2.3), Inches(0.5),
                 title, 18, bold=True, color=DARK)
    add_text_box(slide, left + Inches(0.3), Inches(4.3), Inches(2.3), Inches(2.0),
                 desc, 13, color=GRAY)

    # Arrow
    if i < 3:
        add_text_box(slide, left + Inches(2.7), Inches(4.2), Inches(0.5), Inches(0.5),
                     "→", 24, color=RGBColor(0xCC, 0xCC, 0xCC), alignment=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════
# SLIDE 10: Vision
# ═══════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_shape_bg(slide, DARK)

add_text_box(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(0.5),
             "Workside", 20, bold=True, color=BLUE)

add_text_box(slide, Inches(1.5), Inches(2.5), Inches(10), Inches(1.5),
             "일하는 방식의 다름을\n이해하는 것에서\n성장이 시작됩니다", 44, bold=True, color=WHITE)

add_accent_bar(slide, Inches(1.5), Inches(4.8), Inches(1.0))

add_text_box(slide, Inches(1.5), Inches(5.2), Inches(10), Inches(0.6),
             "Workside는 직장인의 성장과 커리어를\n데이터 기반으로 지원하는 커리어 포탈을 만들어갑니다", 18, color=RGBColor(0xAA, 0xAA, 0xAA))

add_text_box(slide, Inches(1.5), Inches(6.3), Inches(10), Inches(0.5),
             "workside.day", 16, bold=True, color=BLUE)

# ═══════════════════════════════════════════════════
# Save
# ═══════════════════════════════════════════════════
output_path = "/Users/maegbugpeulo/worksider/Workside_서비스소개.pptx"
prs.save(output_path)
print(f"PPT saved: {output_path}")
